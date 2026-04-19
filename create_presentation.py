#!/usr/bin/env python3
"""Генерация презентации для Passport OCR System."""

import sys
sys.path.insert(0, '/Users/arttr/.kilo/skills/pptx')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_PATH = "/Users/arttr/gazprom_hachathon/presentation.pptx"

COLORS = {
    'primary': RGBColor(6, 90, 130),       # #065A82 deep blue
    'secondary': RGBColor(28, 114, 147),  # #1C7293 teal
    'accent': RGBColor(33, 41, 92),        # #21295C midnight
    'white': RGBColor(255, 255, 255),
    'light': RGBColor(242, 242, 242),
    'gray': RGBColor(128, 128, 128),
}

def add_dark_background(slide):
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(5.625)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['accent']
    background.line.fill.background()

def add_slide_title(slide, text):
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

def add_content_box(slide, top, height, text_lines):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), top,
        Inches(9), height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['primary']
    box.line.color.rgb = COLORS['secondary']
    
    tb = box.text_frame
    tb.margin_top = Inches(0.15)
    tb.margin_left = Inches(0.3)
    
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tb.paragraphs[0]
        else:
            p = tb.add_paragraph()
        p.text = line
        p.font.size = Pt(18) if i == 0 else Pt(14)
        p.font.bold = i == 0
        p.font.color.rgb = COLORS['white']
        p.space_after = Pt(8)

def add_stat_callout(slide, left, top, number, label):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top,
        Inches(2.8), Inches(1.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['primary']
    box.line.color.rgb = COLORS['secondary']
    
    tb = box.text_frame
    tb.vertical_alignment = MSO_ANCHOR.MIDDLE
    
    p = tb.paragraphs[0]
    p.text = number
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tb.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLORS['gray']
    p2.alignment = PP_ALIGN.CENTER

def add_list_item(slide, top, index, text):
    left = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.5), top,
        Inches(0.4), Inches(0.4)
    )
    left.fill.solid()
    left.fill.fore_color.rgb = COLORS['secondary']
    left.line.fill.background()
    
    tb = left.text_frame
    p = tb.paragraphs[0]
    p.text = str(index)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    text_box = slide.shapes.add_textbox(Inches(1.1), top, Inches(8.5), Inches(0.5))
    tf = text_box.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = COLORS['white']

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # SLIDE 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Passport OCR System"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.6))
    tf = subtitle.text_frame
    tf.paragraphs[0].text = "Интеллектуальный преобразователь оцифровки паспортных данных"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = COLORS['gray']
    
    team = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.5))
    tf = team.text_frame
    tf.paragraphs[0].text = "Команда: Максим Садреев, Павел Гуров, Артем Трушин"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = COLORS['gray']
    
    hackathon = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.4))
    tf = hackathon.text_frame
    tf.paragraphs[0].text = "Хакатон РТУ МИРЭА 2026"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLORS['secondary']
    
    # SLIDE 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_slide_title(slide, "Проблема")
    
    problems = [
        "Паспорта оборудования поступают на бумажных носителях",
        "Различные формы: стандартные бланки, формы производителя, произвольная форма",
        "Ручной ввод данных — медленно и подвержено ошибкам",
        "Сложность интеграции в систему 1С",
        "Невозможность автоматического поиска и проверки данных",
    ]
    
    for i, prob in enumerate(problems):
        add_list_item(slide, Inches(1.2 + i * 0.8), i + 1, prob)
    
    # SLIDE 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_slide_title(slide, "Решение")
    
    add_content_box(slide, Inches(1.2), Inches(2.0), [
        "Автоматическая оцифровка PDF-документов",
        "OCR с Tesseract / PaddleOCR",
        "LLM-структуризация через Ollama",
    ])
    
    add_content_box(slide, Inches(3.5), Inches(1.8), [
        "Генерация штрих-кодов (Code 39, Code 128, EAN-13)",
        "Экспорт в XLSX для 1С",
        "Веб-интерфейс для загрузки и просмотра",
    ])
    
    # SLIDE 4: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_slide_title(slide, "Архитектура")
    
    add_content_box(slide, Inches(1.2), Inches(3.8), [
        "PDF -> PDFProcessor -> Сегментация страниц (текст/скан)",
        "OCR -> Распознавание текста с confidence",
        "LLMExtractor -> Структуризация через LLM",
        "Генерация штрих-кодов и XLSX-экспорт",
        "Веб-интерфейс: React + FastAPI",
    ])
    
    # SLIDE 5: Features
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_slide_title(slide, "Возможности системы")
    
    features = [
        "Определение типа страниц (текстовая / скан)",
        "OCR: Tesseract и PaddleOCR с оценкой качества",
        "Извлечение таблиц из PDF",
        "LLM-структуризация (Ollama qwen2.5 + llava)",
        "Vision-fallback для проблемных сегментов",
        "Генерация штрих-кодов",
        "Экспорт в JSON и XLSX",
    ]
    
    for i, feat in enumerate(features):
        add_list_item(slide, Inches(1.2 + i * 0.6), i + 1, feat)
    
    # SLIDE 6: Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_slide_title(slide, "Технологический стек")
    
    add_content_box(slide, Inches(1.2), Inches(1.5), [
        "Frontend: React + TypeScript + Vite",
        "Backend: FastAPI + Uvicorn",
    ])
    
    add_content_box(slide, Inches(3.0), Inches(1.5), [
        "OCR: Tesseract, PaddleOCR",
        "PDF: PyMuPDF, pdf2image",
    ])
    
    add_content_box(slide, Inches(4.8), Inches(1.5), [
        "LLM: Ollama (qwen2.5:3b, llava:7b)",
        "Container: Docker + Compose",
    ])
    
    # SLIDE 7: Criteria
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_slide_title(slide, "Критерии оценки")
    
    add_stat_callout(slide, Inches(0.5), Inches(1.5), "0-4", "Точность распознавания")
    add_stat_callout(slide, Inches(3.5), Inches(1.5), "0-2", "Полнота извлеченной информации")
    add_stat_callout(slide, Inches(6.5), Inches(1.5), "0-2", "Удобство интерфейса")
    
    add_stat_callout(slide, Inches(0.5), Inches(3.3), "0-2", "Процент ошибок")
    add_stat_callout(slide, Inches(3.5), Inches(3.3), "+1", "Работа в 1С")
    add_stat_callout(slide, Inches(6.5), Inches(3.3), "+1", "Чек-листы и интеграция")
    
    # SLIDE 8: Team
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_slide_title(slide, "Команда")
    
    team_members = [
        "Максим Садреев",
        "Павел Гуров", 
        "Артем Трушин",
    ]
    
    for i, member in enumerate(team_members):
        name_box = slide.shapes.add_textbox(Inches(1 + i * 3), Inches(2.5), Inches(2.5), Inches(0.6))
        tf = name_box.text_frame
        tf.paragraphs[0].text = member
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS['white']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.5))
    tf = subtitle.text_frame
    tf.paragraphs[0].text = "Хакатон РТУ МИРЭА 2026"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = COLORS['secondary']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    prs.save(OUTPUT_PATH)
    print(f"Презентация сохранена: {OUTPUT_PATH}")

if __name__ == "__main__":
    create_presentation()